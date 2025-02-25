import tkinter as tk
from tkinter import filedialog
from tkinter import ttk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from pathlib import Path
import math


def create_ppt_photo_wall(photo_files, names, output_pptx_path, center_photo_index=0):
    # 创建一个新的 PPT 演示文稿
    prs = Presentation()
    # 选择空白布局，索引为 6
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 确保照片数量和姓名数量一致
    assert len(photo_files) == len(names), "照片数量和姓名数量必须一致"

    # 计算每个图片的宽度和高度
    photo_width = Inches(1)
    photo_height = Inches(1)
    name_height = Inches(0.25)  # 适当减小姓名框高度

    # 幻灯片中心位置
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    center_x = slide_width / 2
    center_y = slide_height / 2

    # 定义统一的字体样式
    font_size = Pt(9)
    font_color = RGBColor(0, 0, 0)  # 黑色
    font_bold = False
    font_italic = False

    # 定义文本框样式
    textbox_fill_color = RGBColor(155, 187, 89)  # 橄榄色填充，强调颜色 3
    textbox_border_color = RGBColor(200, 200, 200)  # 浅色轮廓
    textbox_border_width = Pt(1)

    def set_textbox_style(txBox):
        """设置文本框的样式"""
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = textbox_fill_color

        line = txBox.line
        line.color.rgb = textbox_border_color
        line.width = textbox_border_width

    def set_text_style(tf):
        """设置文本的样式"""
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.color.rgb = font_color
                run.font.bold = font_bold
                run.font.italic = font_italic
            # 设置段落居中对齐
            paragraph.alignment = 1

    # 添加圆心图片
    center_photo = photo_files[center_photo_index]
    center_name = names[center_photo_index]
    center_pic = slide.shapes.add_picture(str(center_photo), center_x - photo_width / 2, center_y - photo_height / 2,
                                          width=photo_width, height=photo_height)
    center_name_left = center_x - photo_width / 2
    center_name_top = center_y + photo_height / 2
    center_txBox = slide.shapes.add_textbox(center_name_left, center_name_top, photo_width, name_height)
    set_textbox_style(center_txBox)
    center_tf = center_txBox.text_frame
    center_tf.text = center_name
    set_text_style(center_tf)

    # 移除圆心图片和姓名
    remaining_photo_files = photo_files[:center_photo_index] + photo_files[center_photo_index + 1:]
    remaining_names = names[:center_photo_index] + names[center_photo_index + 1:]

    # 圆的半径
    radius = Inches(1.5)

    # 图片数量
    num_photos = len(remaining_photo_files)

    # 遍历剩余照片并将它们添加到 PPT 幻灯片上
    for i in range(num_photos):
        # 计算每个图片的角度
        angle = 2 * math.pi * i / num_photos

        # 计算图片的位置
        left = center_x + radius * math.cos(angle) - photo_width / 2
        top = center_y + radius * math.sin(angle) - photo_height / 2

        # 添加图片到幻灯片
        pic = slide.shapes.add_picture(str(remaining_photo_files[i]), left, top, width=photo_width, height=photo_height)

        # 在图片下方添加姓名
        name_left = left
        name_top = top + photo_height
        txBox = slide.shapes.add_textbox(name_left, name_top, photo_width, name_height)
        set_textbox_style(txBox)
        tf = txBox.text_frame
        tf.text = remaining_names[i]
        set_text_style(tf)

    # 保存 PPT
    prs.save(output_pptx_path)


def select_photos_dir():
    photos_dir = filedialog.askdirectory()
    if photos_dir:
        photos_dir_entry.delete(0, tk.END)
        photos_dir_entry.insert(0, photos_dir)


def select_output_pptx_path():
    output_pptx_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                    filetypes=[("PowerPoint Files", "*.pptx")])
    if output_pptx_path:
        output_pptx_path_entry.delete(0, tk.END)
        output_pptx_path_entry.insert(0, output_pptx_path)


def generate_ppt():
    photos_dir = photos_dir_entry.get()
    output_pptx_path = output_pptx_path_entry.get()
    if photos_dir and output_pptx_path:
        photos_dir = Path(photos_dir)
        # 获取所有照片文件
        photo_files = [f for f in photos_dir.iterdir() if f.suffix.lower() in ('.png', '.jpg', '.jpeg')]
        names = [p.name.removesuffix(p.suffix) for p in photo_files]
        try:
            create_ppt_photo_wall(photo_files, names, output_pptx_path, center_photo_index=0)
            status_label.config(text="PPT 生成成功！")
        except Exception as e:
            status_label.config(text=f"生成 PPT 时出错: {str(e)}")
    else:
        status_label.config(text="请选择照片目录和输出 PPT 路径。")


# 创建主窗口
root = tk.Tk()
root.title("照片墙 PPT 生成器")

# 创建样式
style = ttk.Style()

# 创建标签和输入框
photos_dir_label = ttk.Label(root, text="照片目录:")
photos_dir_label.grid(row=0, column=0, padx=10, pady=5, sticky=tk.W)
photos_dir_entry = ttk.Entry(root, width=50)
photos_dir_entry.grid(row=0, column=1, padx=10, pady=5)
photos_dir_button = ttk.Button(root, text="选择目录", command=select_photos_dir)
photos_dir_button.grid(row=0, column=2, padx=10, pady=5)

output_pptx_path_label = ttk.Label(root, text="输出 PPT 路径:")
output_pptx_path_label.grid(row=1, column=0, padx=10, pady=5, sticky=tk.W)
output_pptx_path_entry = ttk.Entry(root, width=50)
output_pptx_path_entry.grid(row=1, column=1, padx=10, pady=5)
output_pptx_path_button = ttk.Button(root, text="选择路径", command=select_output_pptx_path)
output_pptx_path_button.grid(row=1, column=2, padx=10, pady=5)

# 创建生成按钮
generate_button = ttk.Button(root, text="生成 PPT", command=generate_ppt)
generate_button.grid(row=2, column=1, padx=10, pady=20)

# 创建状态标签
status_label = ttk.Label(root, text="")
status_label.grid(row=3, column=1, padx=10, pady=5)

# 运行主循环
root.mainloop()

